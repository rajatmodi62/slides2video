#!/usr/bin/env python3
"""
slides_to_video.py

Takes a .pptx with speaker notes on each slide and produces an .mp4 where
each slide is shown exactly as long as its narration takes, then advances
to the next slide.

Usage:
    python slides_to_video.py my_deck.pptx --out my_deck.mp4
    python slides_to_video.py my_deck.pptx --out my_deck.mp4 --voice my_voice_sample.wav

Requires (system): libreoffice, poppler-utils (pdftoppm), ffmpeg
Requires (python): python-pptx, torch, torchaudio, chatterbox-tts
"""

import argparse
import re
import subprocess
from pathlib import Path

import torch
import torchaudio
from pptx import Presentation
from chatterbox.tts import ChatterboxTTS


def extract_notes(pptx_path):
    prs = Presentation(pptx_path)
    notes = []
    for slide in prs.slides:
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            text = slide.notes_slide.notes_text_frame.text.strip()
        else:
            text = ""
        notes.append(text)
    return notes


def convert_pptx_to_images(pptx_path, out_dir, dpi=150):
    out_dir = Path(out_dir)
    out_dir.mkdir(parents=True, exist_ok=True)

    # Clear any leftover images/PDF from a previous run in this directory -
    # otherwise stale slide-*.png files can get counted alongside the fresh
    # ones and throw off the image/notes count.
    for old in out_dir.glob("slide-*.png"):
        old.unlink()
    for old in out_dir.glob("*.pdf"):
        old.unlink()

    subprocess.run(
        ["libreoffice", "--headless", "--convert-to", "pdf",
         "--outdir", str(out_dir), str(pptx_path)],
        check=True,
    )
    pdf_path = out_dir / (Path(pptx_path).stem + ".pdf")

    prefix = out_dir / "slide"
    subprocess.run(
        ["pdftoppm", "-png", "-r", str(dpi), str(pdf_path), str(prefix)],
        check=True,
    )

    images = sorted(out_dir.glob("slide-*.png"), key=lambda p: int(re.search(r"(\d+)", p.stem).group(1)))
    if not images:
        raise RuntimeError(f"No slide images produced in {out_dir}. Check libreoffice/pdftoppm output above.")
    return images


def split_into_chunks(text, max_chars=280):
    """Chatterbox (like most TTS models) has a max generation length per call
    and will silently truncate long text instead of erroring. Splitting on
    sentence boundaries and generating each chunk separately, then
    concatenating the audio, avoids that truncation."""
    sentences = re.split(r"(?<=[.!?])\s+", text.strip())
    chunks = []
    current = ""
    for sent in sentences:
        if not sent:
            continue
        candidate = (current + " " + sent).strip() if current else sent
        if len(candidate) <= max_chars:
            current = candidate
        else:
            if current:
                chunks.append(current)
            current = sent
    if current:
        chunks.append(current)
    return chunks


def synthesize(model, text, out_path, voice_prompt=None, silence_seconds=1.2):
    text = text.strip()
    if not text:
        # No notes on this slide: hold it briefly in silence rather than crash.
        sr = model.sr
        silence = torch.zeros(1, int(sr * silence_seconds))
        torchaudio.save(str(out_path), silence, sr)
        return

    chunks = split_into_chunks(text)
    wavs = []
    for chunk in chunks:
        if voice_prompt:
            wav = model.generate(chunk, audio_prompt_path=voice_prompt)
        else:
            wav = model.generate(chunk)
        wavs.append(wav)

    combined = wavs[0] if len(wavs) == 1 else torch.cat(wavs, dim=-1)
    torchaudio.save(str(out_path), combined, model.sr)


def make_slide_clip(image_path, audio_path, out_path):
    # -loop 1 holds the still image; -shortest cuts the video to match the
    # audio's length, which is exactly the "hold slide until narration ends" behavior.
    cmd = [
        "ffmpeg", "-y",
        "-loop", "1", "-i", str(image_path),
        "-i", str(audio_path),
        "-vf", "scale=trunc(iw/2)*2:trunc(ih/2)*2",
        "-c:v", "libx264", "-tune", "stillimage",
        "-c:a", "aac", "-b:a", "192k",
        "-pix_fmt", "yuv420p",
        "-shortest",
        str(out_path),
    ]
    result = subprocess.run(cmd, capture_output=True, text=True)
    if result.returncode != 0:
        print(result.stderr)
        raise subprocess.CalledProcessError(result.returncode, cmd)


def concat_clips(clip_paths, out_path, work_dir):
    list_file = Path(work_dir) / "concat_list.txt"
    with open(list_file, "w") as f:
        for p in clip_paths:
            f.write(f"file '{Path(p).resolve()}'\n")
    cmd = [
        "ffmpeg", "-y", "-f", "concat", "-safe", "0",
        "-i", str(list_file), "-c", "copy", str(out_path),
    ]
    result = subprocess.run(cmd, capture_output=True, text=True)
    if result.returncode != 0:
        print(result.stderr)
        raise subprocess.CalledProcessError(result.returncode, cmd)


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("pptx", help="Path to .pptx file with speaker notes")
    parser.add_argument("--out", default="output.mp4", help="Output video path")
    parser.add_argument("--workdir", default="work", help="Directory for intermediate files")
    parser.add_argument("--voice", default=None, help="Optional 5-10s reference .wav to clone a voice")
    parser.add_argument("--dpi", type=int, default=150, help="Slide render resolution")
    args = parser.parse_args()

    work_dir = Path(args.workdir)
    work_dir.mkdir(parents=True, exist_ok=True)

    print("Extracting speaker notes...")
    notes = extract_notes(args.pptx)
    print(f"Found {len(notes)} slides.")

    print("Converting slides to images (libreoffice -> pdf -> png)...")
    images = convert_pptx_to_images(args.pptx, work_dir / "slides", dpi=args.dpi)
    print(f"Rendered {len(images)} slide images.")

    if len(images) != len(notes):
        print(f"WARNING: {len(images)} slide images but {len(notes)} notes entries. "
              f"Something didn't line up 1:1 - check for hidden slides.")

    print("Loading Chatterbox model on GPU...")
    model = ChatterboxTTS.from_pretrained(device="cuda")

    audio_dir = work_dir / "audio"
    audio_dir.mkdir(exist_ok=True)
    clip_dir = work_dir / "clips"
    clip_dir.mkdir(exist_ok=True)

    clip_paths = []
    n = min(len(images), len(notes))
    for i in range(n):
        img, text = images[i], notes[i]
        print(f"[{i + 1}/{n}] Synthesizing narration ({len(text)} chars)...")
        audio_path = audio_dir / f"slide_{i:03d}.wav"
        synthesize(model, text, audio_path, voice_prompt=args.voice)

        print(f"[{i + 1}/{n}] Rendering slide clip...")
        clip_path = clip_dir / f"slide_{i:03d}.mp4"
        make_slide_clip(img, audio_path, clip_path)
        clip_paths.append(clip_path)

    print("Concatenating all slide clips into the final video...")
    concat_clips(clip_paths, args.out, work_dir)
    print(f"\nDone! Video saved to: {args.out}")


if __name__ == "__main__":
    main()